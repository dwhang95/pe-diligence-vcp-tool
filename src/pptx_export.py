"""
pptx_export.py — Shared PowerPoint export for PE Ops Tool Suite.

Usage:
    buf = export_pptx(content, tab, company_name, template=None)
    st.download_button(..., data=buf.getvalue(), ...)

Args:
    content      : dict of {slide_title: markdown_content} for diligence/vcp tabs,
                   or {section_key: markdown_content} for 100day tab
    tab          : "diligence" | "vcp" | "100day"
    company_name : shown on title slide
    template     : BytesIO of an uploaded .pptx, or None for default styling
"""

from __future__ import annotations

import io
import re
from datetime import date
from typing import NamedTuple

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------

SLIDE_W   = Inches(13.33)
SLIDE_H   = Inches(7.5)

_BAR_H    = Inches(1.05)           # navy title bar height
_BODY_L   = Inches(0.45)
_BODY_T   = Inches(1.15)
_BODY_W   = Inches(12.43)
_BODY_H   = Inches(5.85)
_FOOT_T   = Inches(7.1)
_FOOT_H   = Inches(0.32)

MAX_WORDS  = 400
MAX_COLS   = 6
MAX_CELL   = 120

# ---------------------------------------------------------------------------
# Default color palette
# ---------------------------------------------------------------------------

_NAVY    = RGBColor(0x1f, 0x49, 0x7d)   # title bar / table header
_WHITE   = RGBColor(0xff, 0xff, 0xff)
_BODY    = RGBColor(0x1a, 0x1a, 0x1a)
_GRAY    = RGBColor(0x80, 0x80, 0x80)
_ALT     = RGBColor(0xf2, 0xf2, 0xf2)   # alternating table row

# ---------------------------------------------------------------------------
# Per-tab metadata
# ---------------------------------------------------------------------------

_TAB_SUBTITLE = {
    "diligence": "Operational Due Diligence Brief",
    "vcp":       "Value Creation Plan",
    "100day":    "100-Day Operational Plan",
}

# Human-readable titles for 100day section keys
_100DAY_TITLES = {
    "workstreams":       "Workstreams",
    "resource_plan":     "Resource Plan",
    "csuite_assessment": "C-Suite Assessment",
    "org_chart":         "Org Chart — 18-Month Target",
    "ebitda_bridge":     "EBITDA Bridge",
}

# ---------------------------------------------------------------------------
# Style container
# ---------------------------------------------------------------------------

class _Style(NamedTuple):
    font:     str
    title_bg: RGBColor
    title_fg: RGBColor
    body_fg:  RGBColor
    alt_row:  RGBColor
    title_pt: int
    body_pt:  int


_DEFAULT_STYLE = _Style(
    font     = "Arial",
    title_bg = _NAVY,
    title_fg = _WHITE,
    body_fg  = _BODY,
    alt_row  = _ALT,
    title_pt = 24,
    body_pt  = 14,
)

# ---------------------------------------------------------------------------
# Markdown parsing helpers
# ---------------------------------------------------------------------------

def _word_count(text: str) -> int:
    return len(text.split())


def _strip_inline(text: str) -> str:
    """Strip ** / * / ` markdown for plain-text rendering."""
    text = re.sub(r"\*{1,3}(.+?)\*{1,3}", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text.strip()


def _parse_md_table(text: str) -> tuple[str, list[list[str]], str] | None:
    """
    Find the first markdown table in text.
    Returns (text_before, rows, text_after) or None if no table found.
    rows is a list of lists of cell strings (header row first, separator omitted).
    """
    lines = text.splitlines()
    start_idx = end_idx = -1

    for i, line in enumerate(lines):
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            if start_idx == -1:
                start_idx = i
        else:
            if start_idx != -1 and end_idx == -1:
                end_idx = i
                break

    if start_idx == -1:
        return None
    if end_idx == -1:
        end_idx = len(lines)

    table_lines = lines[start_idx:end_idx]
    rows: list[list[str]] = []
    for ln in table_lines:
        stripped = ln.strip()
        if re.match(r"^\|[-:| ]+\|$", stripped):
            continue  # separator row
        cells = [c.strip()[:MAX_CELL] for c in stripped.strip("|").split("|")]
        rows.append(cells)

    if len(rows) < 2:
        return None

    before = "\n".join(lines[:start_idx]).strip()
    after  = "\n".join(lines[end_idx:]).strip()
    return before, rows, after


def _split_by_words(text: str, max_words: int = MAX_WORDS) -> list[str]:
    """
    Split text into chunks of at most max_words words, breaking at paragraph
    (double-newline) boundaries where possible.
    """
    paragraphs = re.split(r"\n{2,}", text.strip())
    chunks: list[str] = []
    current_parts: list[str] = []
    current_count = 0

    for para in paragraphs:
        pw = _word_count(para)
        if current_count + pw > max_words and current_parts:
            chunks.append("\n\n".join(current_parts))
            current_parts = [para]
            current_count = pw
        else:
            current_parts.append(para)
            current_count += pw

    if current_parts:
        chunks.append("\n\n".join(current_parts))

    return chunks or [text]


# ---------------------------------------------------------------------------
# XML cell fill helper
# ---------------------------------------------------------------------------

def _set_cell_fill(cell, color: RGBColor) -> None:
    from lxml import etree
    tcPr = cell._tc.get_or_add_tcPr()
    for existing in tcPr.findall(qn("a:solidFill")):
        tcPr.remove(existing)
    solidFill = etree.SubElement(tcPr, qn("a:solidFill"))
    srgbClr   = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


# ---------------------------------------------------------------------------
# Template helpers
# ---------------------------------------------------------------------------

def _clear_slides(prs: Presentation) -> None:
    """Remove all slides from a presentation while keeping master and layouts."""
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        try:
            prs.part.drop_rel(rId)
        except Exception:
            pass
        sldIdLst.remove(sldId)


def _find_layout(prs: Presentation, *name_hints: str) -> object:
    """Return the first layout whose name matches any hint, or layout index 1."""
    for layout in prs.slide_layouts:
        lname = layout.name.lower()
        if any(h.lower() in lname for h in name_hints):
            return layout
    # Fallback: index 1 (usually "Title and Content"), clamped
    idx = min(1, len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


def _detect_template_style(prs: Presentation) -> _Style:
    """
    Detect dominant font name and accent color from a template's slide master.
    Falls back to _DEFAULT_STYLE values for anything that can't be read.
    """
    font_name = "Arial"
    title_bg  = _NAVY
    title_fg  = _WHITE
    body_fg   = _BODY

    try:
        master = prs.slide_master
        # Font: scan title placeholder in first layout
        for layout in master.slide_layouts[:3]:
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:  # title
                    if ph.has_text_frame:
                        for para in ph.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.name:
                                    font_name = run.font.name
                                    raise StopIteration
    except StopIteration:
        pass
    except Exception:
        pass

    try:
        # Try to read dk1 / accent1 from theme XML
        theme_elem = prs.slide_master.slide_layout_master.element \
            if hasattr(prs.slide_master, "slide_layout_master") else None
        if theme_elem is not None:
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            dk1 = theme_elem.find(".//a:dk1//a:srgbClr", ns)
            if dk1 is not None:
                val = dk1.get("val", "")
                if len(val) == 6:
                    body_fg = RGBColor(
                        int(val[0:2], 16),
                        int(val[2:4], 16),
                        int(val[4:6], 16),
                    )
    except Exception:
        pass

    return _Style(
        font     = font_name,
        title_bg = title_bg,
        title_fg = title_fg,
        body_fg  = body_fg,
        alt_row  = _ALT,
        title_pt = 24,
        body_pt  = 14,
    )


# ---------------------------------------------------------------------------
# Default-mode slide builders
# ---------------------------------------------------------------------------

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_layout(prs: Presentation) -> object:
    """Return the blank layout (index 6 if available, else last)."""
    idx = min(6, len(prs.slide_layouts) - 1)
    return prs.slide_layouts[idx]


def _add_footer(slide, text: str, slide_num: int, total: int, style: _Style) -> None:
    """Add Confidential left + slide number right at slide bottom."""
    W = SLIDE_W
    # Left: Confidential
    tb_l = slide.shapes.add_textbox(_BODY_L, _FOOT_T, Inches(6), _FOOT_H)
    r_l  = tb_l.text_frame.paragraphs[0].add_run()
    r_l.text = text
    r_l.font.size  = Pt(8)
    r_l.font.color.rgb = _GRAY
    r_l.font.name  = style.font
    # Right: slide number
    tb_r = slide.shapes.add_textbox(Inches(11.83), _FOOT_T, Inches(1.5), _FOOT_H)
    p_r  = tb_r.text_frame.paragraphs[0]
    p_r.alignment = PP_ALIGN.RIGHT
    r_r  = p_r.add_run()
    r_r.text = f"{slide_num} / {total}"
    r_r.font.size  = Pt(8)
    r_r.font.color.rgb = _GRAY
    r_r.font.name  = style.font


def _default_title_slide(prs: Presentation, company: str, subtitle: str,
                          today: str, style: _Style) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = style.title_bg

    W, H = SLIDE_W, SLIDE_H

    # Gold rule at vertical midpoint
    rule = slide.shapes.add_shape(1, 0, int(H * 0.52), W, Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(0xc9, 0xa8, 0x4c)
    rule.line.fill.background()

    def _tb(left, top, w, h):
        return slide.shapes.add_textbox(left, top, w, h)

    # Company name
    tb = _tb(Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    run = tf.paragraphs[0].add_run()
    run.text      = company
    run.font.bold = True
    run.font.size = Pt(40)
    run.font.color.rgb = style.title_fg
    run.font.name      = style.font

    # Subtitle (tab type)
    tb2  = _tb(Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.8))
    run2 = tb2.text_frame.paragraphs[0].add_run()
    run2.text      = subtitle
    run2.font.size = Pt(20)
    run2.font.color.rgb = RGBColor(0xc9, 0xa8, 0x4c)
    run2.font.name      = style.font

    # Date
    tb3  = _tb(Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.4))
    run3 = tb3.text_frame.paragraphs[0].add_run()
    run3.text      = today
    run3.font.size = Pt(12)
    run3.font.color.rgb = RGBColor(0xb0, 0xb5, 0xc9)
    run3.font.name      = style.font

    # Footer
    tb4  = _tb(Inches(0.8), Inches(7.1), Inches(11.7), Inches(0.32))
    run4 = tb4.text_frame.paragraphs[0].add_run()
    run4.text      = "Confidential  ·  PE Ops Tool Suite  ·  Pre-diligence draft"
    run4.font.size = Pt(8)
    run4.font.color.rgb = RGBColor(0x6b, 0x70, 0x90)
    run4.font.name      = style.font


def _default_toc_slide(prs: Presentation, titles: list[str],
                       slide_num: int, total: int, style: _Style) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _WHITE

    # Title bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, _BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.title_bg
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(_BODY_L, Inches(0.17), Inches(12.0), Inches(0.7))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text      = "Contents"
    run.font.bold = True
    run.font.size = Pt(style.title_pt)
    run.font.color.rgb = style.title_fg
    run.font.name      = style.font

    # Bullets
    tb2 = slide.shapes.add_textbox(_BODY_L, _BODY_T, _BODY_W, _BODY_H)
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    for i, title in enumerate(titles):
        para = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        run  = para.add_run()
        run.text      = f"  {i + 1}.  {title}"
        run.font.size = Pt(style.body_pt)
        run.font.color.rgb = style.body_fg
        run.font.name      = style.font
        para.space_after = Pt(4)

    _add_footer(slide, "Confidential", slide_num, total, style)


def _default_text_slide(prs: Presentation, title: str, content: str,
                        slide_num: int, total: int, style: _Style) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _WHITE

    # Title bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, _BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.title_bg
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(_BODY_L, Inches(0.17), Inches(12.0), Inches(0.7))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text      = title
    run.font.bold = True
    run.font.size = Pt(style.title_pt)
    run.font.color.rgb = style.title_fg
    run.font.name      = style.font

    # Body text
    tb2 = slide.shapes.add_textbox(_BODY_L, _BODY_T, _BODY_W, _BODY_H)
    tf2 = tb2.text_frame
    tf2.word_wrap = True

    lines = [ln for ln in content.splitlines() if ln.strip()]
    for i, line in enumerate(lines):
        para   = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        clean  = _strip_inline(line.strip())
        is_l1  = clean.startswith(("- ", "* ", "• "))
        is_l2  = clean.startswith(("  - ", "    - "))
        is_hdr = (
            not is_l1 and not is_l2
            and re.match(r"^#{1,3}\s", line.strip())
        )

        run = para.add_run()
        if is_l1:
            run.text    = "  •  " + re.sub(r"^[-*•]\s+", "", clean)
            run.font.size = Pt(style.body_pt)
            para.level  = 0
        elif is_l2:
            run.text    = "      –  " + re.sub(r"^\s+[-*•]\s+", "", clean)
            run.font.size = Pt(max(style.body_pt - 2, 10))
            para.level  = 1
        elif is_hdr:
            clean_hdr   = re.sub(r"^#{1,3}\s+", "", clean)
            run.text    = clean_hdr
            run.font.bold = True
            run.font.size = Pt(style.body_pt)
        else:
            run.text    = clean
            run.font.size = Pt(style.body_pt)

        run.font.color.rgb = style.body_fg
        run.font.name      = style.font
        para.space_after   = Pt(3)

    _add_footer(slide, "Confidential", slide_num, total, style)


def _default_table_slide(prs: Presentation, title: str, rows: list[list[str]],
                         slide_num: int, total: int, style: _Style) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _WHITE

    # Title bar
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, _BAR_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = style.title_bg
    bar.line.fill.background()

    tb = slide.shapes.add_textbox(_BODY_L, Inches(0.17), Inches(12.0), Inches(0.7))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text      = title
    run.font.bold = True
    run.font.size = Pt(style.title_pt)
    run.font.color.rgb = style.title_fg
    run.font.name      = style.font

    n_rows = len(rows)
    n_cols = min(max(len(r) for r in rows), MAX_COLS)

    tbl = slide.shapes.add_table(
        n_rows, n_cols,
        _BODY_L, _BODY_T, _BODY_W, Inches(5.5),
    ).table

    for c in range(n_cols):
        tbl.columns[c].width = int(_BODY_W / n_cols)

    for r_idx, row_data in enumerate(rows):
        is_header = (r_idx == 0)
        fill_color = style.title_bg if is_header else (
            _WHITE if r_idx % 2 == 1 else style.alt_row
        )
        for c_idx in range(n_cols):
            cell = tbl.cell(r_idx, c_idx)
            text = row_data[c_idx] if c_idx < len(row_data) else ""
            cell.text = ""
            tf   = cell.text_frame
            tf.word_wrap = True
            para = tf.paragraphs[0]
            para.space_before = Pt(3)
            para.space_after  = Pt(3)
            run  = para.add_run()
            run.text      = _strip_inline(text)
            run.font.name = style.font
            run.font.bold = is_header
            run.font.size = Pt(11) if is_header else Pt(10)
            run.font.color.rgb = style.title_fg if is_header else style.body_fg
            _set_cell_fill(cell, fill_color)

    _add_footer(slide, "Confidential", slide_num, total, style)


def _default_closing_slide(prs: Presentation, today: str,
                           slide_num: int, style: _Style) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = style.title_bg

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.0))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text      = "Prepared by PE Ops Tool Suite"
    run.font.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = style.title_fg
    run.font.name      = style.font

    tb2  = slide.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.5))
    run2 = tb2.text_frame.paragraphs[0].add_run()
    run2.text      = f"{today}  ·  Confidential  ·  Pre-diligence draft"
    run2.font.size = Pt(13)
    run2.font.color.rgb = RGBColor(0xb0, 0xb5, 0xc9)
    run2.font.name      = style.font


# ---------------------------------------------------------------------------
# Template-mode slide builders
# ---------------------------------------------------------------------------

def _get_placeholder(slide, idx: int):
    """Return a placeholder by idx, or None."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _set_ph_text(ph, text: str, bold: bool = False,
                 font_name: str = "Arial", pt: int = 24,
                 color: RGBColor | None = None) -> None:
    """Set plain text in a placeholder, clearing existing content."""
    tf = ph.text_frame
    tf.clear()
    para = tf.paragraphs[0]
    run  = para.add_run()
    run.text      = text
    run.font.bold = bold
    run.font.size = Pt(pt)
    run.font.name = font_name
    if color:
        run.font.color.rgb = color


def _template_title_slide(prs: Presentation, layout, company: str,
                          subtitle: str, today: str, style: _Style) -> None:
    slide = prs.slides.add_slide(layout)
    title_ph = _get_placeholder(slide, 0)
    body_ph  = _get_placeholder(slide, 1)

    if title_ph:
        _set_ph_text(title_ph, company, bold=True,
                     font_name=style.font, pt=style.title_pt + 8)
    else:
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.4))
        r  = tb.text_frame.paragraphs[0].add_run()
        r.text = company; r.font.bold = True
        r.font.size = Pt(36); r.font.name = style.font

    if body_ph:
        _set_ph_text(body_ph, f"{subtitle}\n{today}",
                     font_name=style.font, pt=style.body_pt)
    else:
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.8))
        r2  = tb2.text_frame.paragraphs[0].add_run()
        r2.text = subtitle; r2.font.size = Pt(18); r2.font.name = style.font


def _template_text_slide(prs: Presentation, layout, title: str, content: str,
                         slide_num: int, total: int, style: _Style) -> None:
    slide    = prs.slides.add_slide(layout)
    title_ph = _get_placeholder(slide, 0)
    body_ph  = _get_placeholder(slide, 1)

    if title_ph:
        _set_ph_text(title_ph, title, bold=True,
                     font_name=style.font, pt=style.title_pt)
    else:
        tb = slide.shapes.add_textbox(_BODY_L, Inches(0.17), Inches(12.0), Inches(0.7))
        r  = tb.text_frame.paragraphs[0].add_run()
        r.text = title; r.font.bold = True
        r.font.size = Pt(style.title_pt); r.font.name = style.font

    # Determine body bounds
    if body_ph:
        body_left  = body_ph.left
        body_top   = body_ph.top
        body_width = body_ph.width
        target_frame = body_ph.text_frame
    else:
        tb2 = slide.shapes.add_textbox(_BODY_L, _BODY_T, _BODY_W, _BODY_H)
        body_left  = _BODY_L
        body_top   = _BODY_T
        body_width = _BODY_W
        target_frame = tb2.text_frame

    target_frame.word_wrap = True
    lines = [ln for ln in content.splitlines() if ln.strip()]
    for i, line in enumerate(lines):
        para  = target_frame.paragraphs[0] if i == 0 else target_frame.add_paragraph()
        clean = _strip_inline(line.strip())
        is_l1 = clean.startswith(("- ", "* ", "• "))
        is_l2 = clean.startswith(("  - ", "    - "))
        run   = para.add_run()
        if is_l1:
            run.text  = "  •  " + re.sub(r"^[-*•]\s+", "", clean)
        elif is_l2:
            run.text  = "      –  " + re.sub(r"^\s+[-*•]\s+", "", clean)
            para.level = 1
        else:
            run.text  = clean
        run.font.size = Pt(style.body_pt)
        run.font.name = style.font
        para.space_after = Pt(3)

    # Add slide number in bottom right as textbox (template footer placeholder may exist)
    _num_tb = slide.shapes.add_textbox(Inches(11.83), _FOOT_T, Inches(1.5), _FOOT_H)
    p_num   = _num_tb.text_frame.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    r_num   = p_num.add_run()
    r_num.text = f"{slide_num} / {total}"
    r_num.font.size = Pt(8); r_num.font.name = style.font


def _template_table_slide(prs: Presentation, layout, title: str,
                          rows: list[list[str]], slide_num: int, total: int,
                          style: _Style) -> None:
    slide    = prs.slides.add_slide(layout)
    title_ph = _get_placeholder(slide, 0)
    body_ph  = _get_placeholder(slide, 1)

    if title_ph:
        _set_ph_text(title_ph, title, bold=True,
                     font_name=style.font, pt=style.title_pt)
    else:
        tb = slide.shapes.add_textbox(_BODY_L, Inches(0.17), Inches(12.0), Inches(0.7))
        r  = tb.text_frame.paragraphs[0].add_run()
        r.text = title; r.font.bold = True
        r.font.size = Pt(style.title_pt); r.font.name = style.font

    # Position table inside body placeholder if available
    if body_ph:
        tbl_l = body_ph.left
        tbl_t = body_ph.top
        tbl_w = body_ph.width
        tbl_h = body_ph.height
        # Hide the body placeholder text
        body_ph.text_frame.clear()
    else:
        tbl_l, tbl_t, tbl_w, tbl_h = _BODY_L, _BODY_T, _BODY_W, Inches(5.5)

    n_rows = len(rows)
    n_cols = min(max(len(r) for r in rows), MAX_COLS)

    tbl = slide.shapes.add_table(n_rows, n_cols, tbl_l, tbl_t, tbl_w, tbl_h).table
    for c in range(n_cols):
        tbl.columns[c].width = int(tbl_w / n_cols)

    for r_idx, row_data in enumerate(rows):
        is_header  = (r_idx == 0)
        fill_color = style.title_bg if is_header else (
            _WHITE if r_idx % 2 == 1 else style.alt_row
        )
        for c_idx in range(n_cols):
            cell = tbl.cell(r_idx, c_idx)
            text = row_data[c_idx] if c_idx < len(row_data) else ""
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            para.space_before = Pt(3)
            para.space_after  = Pt(3)
            run  = para.add_run()
            run.text      = _strip_inline(text)
            run.font.name = style.font
            run.font.bold = is_header
            run.font.size = Pt(11) if is_header else Pt(10)
            run.font.color.rgb = style.title_fg if is_header else style.body_fg
            _set_cell_fill(cell, fill_color)

    r_num_tb = slide.shapes.add_textbox(Inches(11.83), _FOOT_T, Inches(1.5), _FOOT_H)
    p_num    = r_num_tb.text_frame.paragraphs[0]
    p_num.alignment = PP_ALIGN.RIGHT
    rr = p_num.add_run()
    rr.text = f"{slide_num} / {total}"
    rr.font.size = Pt(8); rr.font.name = style.font


# ---------------------------------------------------------------------------
# Section renderer — handles splits, tables, and both modes
# ---------------------------------------------------------------------------

def _render_section(
    prs: Presentation,
    title: str,
    content: str,
    slide_num: int,
    total: int,
    style: _Style,
    template_mode: bool,
    title_layout = None,
    content_layout = None,
) -> int:
    """
    Render one section as 1+ slides.
    Returns the updated current slide number.
    """
    def text_slide(t, c, n):
        if template_mode:
            _template_text_slide(prs, content_layout, t, c, n, total, style)
        else:
            _default_text_slide(prs, t, c, n, total, style)

    def table_slide(t, rows, n):
        if template_mode:
            _template_table_slide(prs, content_layout, t, rows, n, total, style)
        else:
            _default_table_slide(prs, t, rows, n, total, style)

    # Detect table in content
    table_result = _parse_md_table(content)

    if table_result is not None:
        before_text, rows, after_text = table_result

        # Split rows if > MAX_COLS (keep col 0 as anchor, paginate rest)
        n_data_cols = max(len(r) for r in rows)
        if n_data_cols > MAX_COLS:
            col_groups = [list(range(min(MAX_COLS, n_data_cols)))]
            remaining  = list(range(MAX_COLS, n_data_cols))
            while remaining:
                batch = [0] + remaining[:MAX_COLS - 1]
                col_groups.append(batch)
                remaining = remaining[MAX_COLS - 1:]
        else:
            col_groups = [list(range(n_data_cols))]

        # Text before table
        if before_text.strip() and _word_count(before_text) > 10:
            for chunk in _split_by_words(before_text):
                text_slide(title, chunk, slide_num)
                slide_num += 1
                title = title + " (cont'd)"

        # Table slide(s)
        for group_num, col_indices in enumerate(col_groups):
            group_rows = [[row[ci] if ci < len(row) else "" for ci in col_indices]
                          for row in rows]
            tbl_title  = title if group_num == 0 else title + " (cont'd)"
            table_slide(tbl_title, group_rows, slide_num)
            slide_num += 1
            title = title + " (cont'd)"

        # Text after table
        if after_text.strip() and _word_count(after_text) > 10:
            for chunk in _split_by_words(after_text):
                text_slide(title, chunk, slide_num)
                slide_num += 1
                title = title + " (cont'd)"

    else:
        # Pure text — split on word count
        chunks = _split_by_words(content)
        for chunk_idx, chunk in enumerate(chunks):
            t = title if chunk_idx == 0 else title + " (cont'd)"
            text_slide(t, chunk, slide_num)
            slide_num += 1

    return slide_num


# ---------------------------------------------------------------------------
# Master export function
# ---------------------------------------------------------------------------

def export_pptx(
    content: dict[str, str],
    tab: str,
    company_name: str,
    template: io.BytesIO | None = None,
) -> io.BytesIO:
    """
    Build a PowerPoint presentation from generated section content.

    Args:
        content:      {section_key_or_title: markdown_text}
                      For tab="100day" keys are internal (workstreams, etc.).
                      For tab="diligence"/"vcp" keys are human-readable titles.
        tab:          "diligence" | "vcp" | "100day"
        company_name: shown on title slide
        template:     BytesIO of an uploaded .pptx, or None for default styling

    Returns:
        BytesIO positioned at 0 ready for st.download_button
    """
    today    = date.today().strftime("%Y-%m-%d")
    subtitle = _TAB_SUBTITLE.get(tab, "PE Ops Output")

    # Resolve section display titles and filter empty/error content
    if tab == "100day":
        sections: dict[str, str] = {
            _100DAY_TITLES.get(k, k.replace("_", " ").title()): v
            for k, v in content.items()
            if v and not v.startswith("_Error")
        }
    else:
        sections = {k: v for k, v in content.items()
                    if v and not v.startswith("_Error")}

    section_titles = list(sections.keys())

    # ── Template mode ────────────────────────────────────────────────────────
    if template is not None:
        try:
            prs   = Presentation(template)
            style = _detect_template_style(prs)
            _clear_slides(prs)
            title_layout   = prs.slide_layouts[0]  # Title Slide
            content_layout = _find_layout(prs, "Title and Content", "content", "body")
            template_mode  = True
        except Exception:
            # Template failed to load — fall through to default mode
            prs            = _new_prs()
            style          = _DEFAULT_STYLE
            title_layout   = None
            content_layout = None
            template_mode  = False
    else:
        prs            = _new_prs()
        style          = _DEFAULT_STYLE
        title_layout   = None
        content_layout = None
        template_mode  = False

    # Count total slides: 1 title + 1 TOC + content slides + 1 closing
    # Estimate content slides (1 per section, some may split — use 1.5x heuristic)
    estimated_content = max(len(sections), 1)
    total_approx = 2 + estimated_content + 1  # title + toc + content + closing

    # ── Slide 1: Title ───────────────────────────────────────────────────────
    if template_mode:
        _template_title_slide(prs, title_layout, company_name, subtitle, today, style)
    else:
        _default_title_slide(prs, company_name, subtitle, today, style)

    # ── Slide 2: TOC ─────────────────────────────────────────────────────────
    if template_mode:
        _template_text_slide(
            prs, content_layout, "Contents",
            "\n".join(f"- {t}" for t in section_titles),
            2, total_approx, style,
        )
    else:
        _default_toc_slide(prs, section_titles, 2, total_approx, style)

    # ── Content slides ────────────────────────────────────────────────────────
    slide_num = 3
    for sec_title, sec_content in sections.items():
        slide_num = _render_section(
            prs, sec_title, sec_content,
            slide_num, total_approx,
            style, template_mode,
            title_layout=title_layout,
            content_layout=content_layout,
        )

    # ── Closing slide ─────────────────────────────────────────────────────────
    if template_mode:
        _template_text_slide(
            prs, content_layout,
            "Prepared by PE Ops Tool Suite",
            today,
            slide_num, slide_num, style,
        )
    else:
        _default_closing_slide(prs, today, slide_num, style)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
